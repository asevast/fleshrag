import os
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation


def parse_docx(path: str) -> str:
    doc = Document(path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    return "\n".join(paragraphs)


def parse_xlsx(path: str) -> str:
    wb = load_workbook(path, data_only=True)
    lines = []
    for sheet in wb.worksheets:
        lines.append(f"\n--- Лист: {sheet.title} ---")
        for row in sheet.iter_rows(values_only=True):
            row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
            if row_text.strip():
                lines.append(row_text)
    return "\n".join(lines)


def parse_pptx(path: str) -> str:
    prs = Presentation(path)
    lines = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        lines.append(f"\n--- Слайд {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text)
    return "\n".join(lines)


def parse_txt(path: str) -> str:
    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        return f.read()
